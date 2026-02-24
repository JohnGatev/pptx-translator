from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
import json
import re
import time
import zipfile
from typing import Dict, Iterable, List, Sequence, Tuple
import xml.etree.ElementTree as ET

import requests
import streamlit as st
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from rapidfuzz import fuzz


DEFAULT_BASE_URL = "https://ai-research-proxy.azurewebsites.net"
DEFAULT_MODEL = "gpt-oss-120b"
DEFAULT_TEMPERATURE = 0.2
DEFAULT_TIMEOUT = 120
LEXICON_PDF_NAME = "Finance 1 Lexicon ENG NL.pdf"
FUZZY_MATCH_THRESHOLD = 88
OCR_IMAGE_ENABLED = True
LINE_BREAK_TOKEN = "[[[BR]]]"



@dataclass(frozen=True)
class TranslationConfig:
    base_url: str
    model: str
    target_language: str
    temperature: float
    timeout: int
    api_key: str


@dataclass(frozen=True)
class LexiconEntry:
    term: str
    translation: str
    raw: str


@dataclass(frozen=True)
class LexiconStore:
    entries: List[LexiconEntry]
    normalized_terms: List[str]
    term_tokens: List[set[str]]


def normalize_base_url(base_url: str) -> str:
    return base_url.rstrip("/")


def load_finance_lexicon() -> LexiconStore | None:
    lexicon_path = Path(__file__).parent / LEXICON_PDF_NAME
    if not lexicon_path.exists():
        return None
    try:
        from pypdf import PdfReader
    except ImportError:
        return None

    reader = PdfReader(str(lexicon_path))
    lines: List[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        lines.extend(text.splitlines())

    entries: List[LexiconEntry] = []
    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            continue
        if re.fullmatch(r"(english|engels)\s*(dutch|nederlands)", line, re.IGNORECASE):
            continue
        # prefer split on multi-space columns
        if re.search(r"\s{2,}", raw_line):
            parts = re.split(r"\s{2,}", raw_line.strip())
        else:
            parts = re.split(r"\s*(?:=>|->|–|—|-|:|\t)\s*", line, maxsplit=1)
        if len(parts) >= 2:
            english = parts[0].strip()
            dutch = parts[1].strip()
            if english and dutch:
                raw = f"{english} => {dutch}"
                entries.append(LexiconEntry(term=english, translation=dutch, raw=raw))
        else:
            entries.append(LexiconEntry(term=line, translation="", raw=line))

    if not entries:
        return None
    return build_lexicon_store(entries)


def build_lexicon_store(entries: List[LexiconEntry]) -> LexiconStore:
    normalized_terms: List[str] = []
    term_tokens: List[set[str]] = []
    for entry in entries:
        normalized = normalize_phrase(entry.term)
        normalized_terms.append(normalized)
        term_tokens.append(tokens_with_stems(normalized))
    return LexiconStore(entries=entries, normalized_terms=normalized_terms, term_tokens=term_tokens)


def normalize_phrase(text: str) -> str:
    cleaned = re.sub(r"\[\[\[(?:RUN_(?:\d+|END)|BR)\]\]\]", " ", text)
    cleaned = re.sub(r"[^A-Za-zÀ-ÿ0-9]+", " ", cleaned.lower())
    return re.sub(r"\s+", " ", cleaned).strip()


def stem_token(token: str) -> str:
    for suffix in ("ing", "ed", "es", "s"):
        if token.endswith(suffix) and len(token) > len(suffix) + 2:
            return token[: -len(suffix)]
    return token


def tokens_with_stems(normalized_text: str) -> set[str]:
    if not normalized_text:
        return set()
    tokens = normalized_text.split()
    stems = {stem_token(token) for token in tokens}
    return set(tokens) | stems

def encode_line_breaks(text: str) -> str:
    if not text:
        return text
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    return normalized.replace("\n", LINE_BREAK_TOKEN)


def decode_line_breaks(text: str) -> str:
    if not text:
        return text
    return text.replace(LINE_BREAK_TOKEN, "\n")


def fuzzy_token_match(term_tokens: set[str], doc_tokens: Sequence[str], threshold: int) -> bool:
    if not term_tokens or not doc_tokens:
        return False
    filtered = [token for token in term_tokens if len(token) >= 3]
    if not filtered:
        return False
    for token in filtered:
        if token in doc_tokens:
            continue
        best = max(fuzz.ratio(token, candidate) for candidate in doc_tokens)
        if best < threshold:
            return False
    return True


def select_lexicon_matches(text: str, lexicon: LexiconStore) -> List[LexiconEntry]:
    normalized_text = normalize_phrase(text)
    if not normalized_text:
        return []
    doc_tokens_set = tokens_with_stems(normalized_text)
    doc_tokens_list = list(doc_tokens_set) or normalized_text.split()
    matches: List[LexiconEntry] = []
    for idx, entry in enumerate(lexicon.entries):
        term_norm = lexicon.normalized_terms[idx]
        if not term_norm:
            continue
        if term_norm in normalized_text:
            matches.append(entry)
            continue
        term_tokens = lexicon.term_tokens[idx]
        if term_tokens and term_tokens.issubset(doc_tokens_set):
            matches.append(entry)
            continue
        if FUZZY_MATCH_THRESHOLD and fuzzy_token_match(term_tokens, doc_tokens_list, FUZZY_MATCH_THRESHOLD):
            matches.append(entry)
    return dedupe_lexicon_matches(matches)


def dedupe_lexicon_matches(matches: Sequence[LexiconEntry]) -> List[LexiconEntry]:
    seen = set()
    deduped: List[LexiconEntry] = []
    for entry in matches:
        key = (entry.term.lower(), entry.translation.lower())
        if key in seen:
            continue
        seen.add(key)
        deduped.append(entry)
    return deduped


def merge_lexicon_matches(match_groups: Iterable[Sequence[LexiconEntry]]) -> List[LexiconEntry]:
    combined: List[LexiconEntry] = []
    for group in match_groups:
        combined.extend(group)
    return dedupe_lexicon_matches(combined)


def format_lexicon_context(matches: Sequence[LexiconEntry]) -> str:
    if not matches:
        return ""
    lines = []
    for entry in matches:
        if entry.translation:
            lines.append(f"- {entry.term} => {entry.translation}")
        else:
            lines.append(f"- {entry.raw}")
    return "\n".join(lines)


def rag_instructions(rag_context: str | None) -> str:
    if not rag_context:
        return ""
    return (
        "\n\nFinance lexicon context. Use these preferred translations when relevant "
        "(including close variations). Do not output this list:\n"
        f"{rag_context}\n"
    )


def strip_glossary_spillover(text: str) -> str:
    if not text:
        return text
    cleaned_lines = []
    for line in text.splitlines():
        lower = line.lower()
        if "finance lexicon" in lower:
            continue
        if re.match(r"^\s*[-•*]?\s*[^>]+=>\s*[^>]+", line):
            continue
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines).strip()


def strip_wrapping_quotes(text: str) -> str:
    stripped = text.strip()
    if len(stripped) >= 2 and stripped[0] == stripped[-1] and stripped[0] in {"'", "\""}:
        return stripped[1:-1]
    return text


def translation_contains_term(translation: str, required: str) -> bool:
    required_norm = normalize_phrase(required)
    if not required_norm:
        return True
    translation_norm = normalize_phrase(translation)
    if required_norm in translation_norm:
        return True
    required_tokens = tokens_with_stems(required_norm)
    translation_tokens = tokens_with_stems(translation_norm)
    if required_tokens and required_tokens.issubset(translation_tokens):
        return True
    if FUZZY_MATCH_THRESHOLD and translation_tokens:
        return fuzzy_token_match(required_tokens, list(translation_tokens), max(80, FUZZY_MATCH_THRESHOLD - 8))
    return False


def needs_repair(source_text: str, translated_text: str) -> bool:
    source_norm = normalize_phrase(source_text)
    translated_norm = normalize_phrase(translated_text)
    if not source_norm or not translated_norm:
        return False
    if len(translated_norm) < max(12, int(len(source_norm) * 0.6)):
        return True
    source_tokens = source_norm.split()
    translated_tokens = translated_norm.split()
    if source_tokens and translated_tokens:
        ratio = len(translated_tokens) / max(1, len(source_tokens))
        if ratio < 0.65:
            return True
    return False


def repair_translation(
    source_text: str,
    translated_text: str,
    config: TranslationConfig,
    glossary_matches: Sequence[LexiconEntry] | None = None,
) -> str:
    glossary_text = format_lexicon_context(glossary_matches or [])
    glossary_note = (
        "\n\nUse these required translations exactly when they appear in the source:\n"
        f"{glossary_text}\n"
        if glossary_text
        else ""
    )
    system_prompt = (
        "You are repairing a translation. Preserve every word and all placeholders "
        "like [[[RUN_0]]], keep them unchanged and in the same order. Do not omit any text, "
        "numbers, or punctuation. Rewrite the translation so it is complete and faithful, "
        f"and ensure glossary terms are used exactly.{glossary_note}"
    )
    user_prompt = (
        "Original text:\n"
        f"{source_text}\n\n"
        "Current translation:\n"
        f"{translated_text}\n\n"
        "Return only the corrected translation."
    )
    return request_chat_completion(
        [{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
        config,
    )


def postprocess_translation(
    source_text: str,
    translated_text: str,
    config: TranslationConfig,
    glossary_matches: Sequence[LexiconEntry] | None = None,
) -> str:
    cleaned = strip_glossary_spillover(strip_wrapping_quotes(translated_text))
    matches = glossary_matches or []
    missing = [
        entry
        for entry in matches
        if entry.translation and not translation_contains_term(cleaned, entry.translation)
    ]
    if missing or needs_repair(source_text, cleaned):
        repaired = repair_translation(source_text, cleaned, config, missing or matches)
        cleaned = strip_glossary_spillover(strip_wrapping_quotes(repaired))
    return cleaned


def split_whitespace(text: str) -> tuple[str, str, str]:
    match = re.match(r"^(\s*)(.*?)(\s*)$", text, re.DOTALL)
    if not match:
        return "", text, ""
    return match.group(1), match.group(2), match.group(3)


def build_translation_messages(
    text: str,
    config: TranslationConfig,
    rag_context: str | None = None,
) -> List[Dict[str, object]]:
    rag_note = rag_instructions(rag_context)
    return [
        {
            "role": "system",
            "content": (
                "You are a professional translation engine. Translate the user's text "
                f"to {config.target_language}. Preserve meaning, punctuation, and line "
                "breaks. Do not omit any words or sentences; translate verbatim. "
                "The text may include placeholders like [[[RUN_0]]] or [[[BR]]]. "
                "Keep all placeholders unchanged and in the same order. Only translate the "
                "text between placeholders. Do not add commentary or quotes."
                f"{rag_note}"
            ),
        },
        {"role": "user", "content": text},
    ]


def request_chat_completion(messages: List[Dict[str, object]], config: TranslationConfig) -> str:
    base_url = normalize_base_url(config.base_url)
    payload = {
        "model": config.model,
        "messages": messages,
        "temperature": config.temperature,
    }
    headers = {
        "Authorization": f"Bearer {config.api_key}",
        "Content-Type": "application/json",
    }
    response = requests.post(
        f"{base_url}/v1/chat/completions",
        json=payload,
        headers=headers,
        timeout=config.timeout,
    )
    response.raise_for_status()
    data = response.json()
    content = data.get("choices", [{}])[0].get("message", {}).get("content")
    if content is None:
        raise RuntimeError("Translation response did not contain translated text.")
    return content


def translate_text(text: str, config: TranslationConfig, rag_context: str | None = None) -> str:
    return request_chat_completion(build_translation_messages(text, config, rag_context), config)


def cached_translate(text: str, config: TranslationConfig, rag_context: str | None = None) -> str:
    return translate_text(text, config, rag_context)


@dataclass
class TranslationSegment:
    segment_id: str
    text: str


@dataclass
class ParagraphInfo:
    runs: List
    prefixes: List[str]
    cores: List[str]
    suffixes: List[str]
    translate_mask: List[bool]
    shape: object
    original_length: int
    is_table: bool
    slide_index: int
    is_title: bool


def is_numeric_text(text: str) -> bool:
    return bool(re.fullmatch(r"[\d\s.,:%€$£¥+-]+", text.strip()))


def is_math_text(text: str) -> bool:
    if not text:
        return False
    math_symbols = set("=<>±√∑∫∏≈≠≤≥→←⇒⇔∞πσΔΩβμτ")
    if any(symbol in text for symbol in math_symbols):
        return True
    if re.search(r"[A-Za-z]\s*=\s*[\d]", text):
        return True
    if re.search(r"[A-Za-z]\s*[_^]", text):
        return True
    return False


def is_off_canvas(shape, slide_width: int, slide_height: int) -> bool:
    left = getattr(shape, "left", 0)
    top = getattr(shape, "top", 0)
    width = getattr(shape, "width", 0)
    height = getattr(shape, "height", 0)
    if width <= 0 or height <= 0:
        return True
    right = left + width
    bottom = top + height
    return right < 0 or bottom < 0 or left > slide_width or top > slide_height


def should_skip_shape(shape, slide_width: int, slide_height: int, seen_keys: set) -> bool:
    if getattr(shape, "visible", True) is False:
        return True
    if is_off_canvas(shape, slide_width, slide_height):
        return True
    if shape.is_placeholder:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in {
            PP_PLACEHOLDER.DATE,
            PP_PLACEHOLDER.FOOTER,
            PP_PLACEHOLDER.SLIDE_NUMBER,
        }:
            return True
    text = extract_shape_text(shape)
    if not text.strip():
        return False
    left = int(getattr(shape, "left", 0))
    top = int(getattr(shape, "top", 0))
    width = int(getattr(shape, "width", 0))
    height = int(getattr(shape, "height", 0))
    key = (left, top, width, height, text.strip())
    if key in seen_keys:
        return True
    seen_keys.add(key)
    return False


def extract_shape_text(shape) -> str:
    if shape.has_text_frame:
        return shape.text_frame.text or ""
    if shape.has_table:
        parts = []
        for row in shape.table.rows:
            for cell in row.cells:
                parts.append(cell.text or "")
        return "\n".join(parts)
    return ""


def iter_slide_shapes(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_slide_shapes(shape.shapes)
        else:
            yield shape


def iter_paragraph_runs(shape) -> Iterable[Tuple[List, bool]]:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from iter_paragraph_runs(child)
        return

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.runs:
                yield list(paragraph.runs), False

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    if paragraph.runs:
                        yield list(paragraph.runs), True


def collect_paragraphs(presentation: Presentation) -> List[ParagraphInfo]:
    paragraphs: List[ParagraphInfo] = []
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    for slide_index, slide in enumerate(presentation.slides):
        seen_keys: set = set()
        for shape in iter_slide_shapes(slide.shapes):
            if should_skip_shape(shape, slide_width, slide_height, seen_keys):
                continue
            for runs, is_table in iter_paragraph_runs(shape):
                is_title = False
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE}:
                        is_title = True
                prefixes = []
                cores = []
                suffixes = []
                translate_mask = []
                original_length = 0
                for run in runs:
                    prefix, core, suffix = split_whitespace(run.text or "")
                    can_translate = bool(core)
                    if core:
                        font_name = getattr(getattr(run, "font", None), "name", None)
                        if font_name and "math" in font_name.lower():
                            can_translate = False
                        if is_math_text(core) or is_numeric_text(core):
                            can_translate = False
                    prefixes.append(prefix)
                    cores.append(core)
                    suffixes.append(suffix)
                    translate_mask.append(can_translate)
                    if can_translate:
                        original_length += len(core)
                if any(translate_mask):
                    paragraphs.append(
                        ParagraphInfo(
                            runs=runs,
                            prefixes=prefixes,
                            cores=cores,
                            suffixes=suffixes,
                            translate_mask=translate_mask,
                            shape=shape,
                            original_length=original_length,
                            is_table=is_table,
                            slide_index=slide_index,
                            is_title=is_title,
                        )
                    )
    return paragraphs


def iter_docx_paragraph_runs(document: Document) -> Iterable[Tuple[List, bool]]:
    for paragraph in document.paragraphs:
        if paragraph.runs:
            yield list(paragraph.runs), False
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if paragraph.runs:
                        yield list(paragraph.runs), True


def collect_docx_paragraphs(document: Document) -> List[ParagraphInfo]:
    paragraphs: List[ParagraphInfo] = []
    for runs, is_table in iter_docx_paragraph_runs(document):
        prefixes = []
        cores = []
        suffixes = []
        translate_mask = []
        original_length = 0
        for run in runs:
            prefix, core, suffix = split_whitespace(run.text or "")
            can_translate = bool(core)
            if core:
                if is_math_text(core) or is_numeric_text(core):
                    can_translate = False
            prefixes.append(prefix)
            cores.append(core)
            suffixes.append(suffix)
            translate_mask.append(can_translate)
            if can_translate:
                original_length += len(core)
        if any(translate_mask):
            paragraphs.append(
                ParagraphInfo(
                    runs=runs,
                    prefixes=prefixes,
                    cores=cores,
                    suffixes=suffixes,
                    translate_mask=translate_mask,
                    shape=None,
                    original_length=original_length,
                    is_table=True,
                    slide_index=0,
                    is_title=False,
                )
            )
    return paragraphs


def build_segments_json(segments: Sequence[TranslationSegment]) -> str:
    payload = {"segments": [{"id": seg.segment_id, "text": seg.text} for seg in segments]}
    return json.dumps(payload, ensure_ascii=False)


def parse_segments_json(raw: str) -> List[Dict[str, str]]:
    raw = raw.strip()
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", raw, re.DOTALL)
        if not match:
            match = re.search(r"\[.*\]", raw, re.DOTALL)
            if not match:
                raise
            data = json.loads(match.group(0))
        else:
            data = json.loads(match.group(0))
    if isinstance(data, list):
        return data
    segments = data.get("segments")
    if not isinstance(segments, list):
        raise ValueError("Invalid translation response format.")
    return segments


def translate_segments(
    segments: Sequence[TranslationSegment],
    config: TranslationConfig,
    rag_context: str | None = None,
) -> Dict[str, str]:
    rag_note = rag_instructions(rag_context)
    system_prompt = (
        "You are a professional translation engine. Translate the provided JSON segments "
        f"to {config.target_language}. Return ONLY valid JSON with the same 'segments' array. "
        "Each item must have the same 'id' and the translated 'text'. Preserve all placeholders "
        "like [[[RUN_0]]] or [[[BR]]], keep them unchanged and in the same order. "
        "Do not omit any words or sentences; translate verbatim. Do not add commentary."
        f"{rag_note}"
    )
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": build_segments_json(segments)},
    ]
    content = request_chat_completion(messages, config)
    parsed = parse_segments_json(content)
    results: Dict[str, str] = {}
    for item in parsed:
        segment_id = str(item.get("id", ""))
        text = item.get("text")
        if not segment_id or text is None:
            raise ValueError("Translation output missing id/text.")
        results[segment_id] = text
    if len(results) != len(segments):
        raise ValueError("Translation output segment count mismatch.")
    return results


def translate_segments_cached(
    segments: Tuple[Tuple[str, str], ...],
    config: TranslationConfig,
    rag_context: str | None = None,
) -> Dict[str, str]:
    segment_objs = [TranslationSegment(segment_id=seg_id, text=text) for seg_id, text in segments]
    return translate_segments(segment_objs, config, rag_context)


def translate_text_with_postprocess(
    source_text: str,
    config: TranslationConfig,
    glossary: Sequence[LexiconEntry] | None = None,
) -> str:
    prepared = encode_line_breaks(source_text)
    rag_context = format_lexicon_context(glossary) if glossary else None
    translated = cached_translate(prepared, config, rag_context)
    translated = postprocess_translation(prepared, translated, config, glossary)
    return decode_line_breaks(translated)


def batch_segments(segments: Sequence[TranslationSegment], max_chars: int = 4000, max_items: int = 20):
    batch: List[TranslationSegment] = []
    total_chars = 0
    for segment in segments:
        segment_len = len(segment.text)
        if batch and (len(batch) >= max_items or total_chars + segment_len > max_chars):
            yield batch
            batch = []
            total_chars = 0
        batch.append(segment)
        total_chars += segment_len
    if batch:
        yield batch


def batch_limits() -> Tuple[int, int]:
    return 30000, 120


def count_translation_batches(paragraphs: List[ParagraphInfo]) -> int:
    grouped: Dict[int, List[ParagraphInfo]] = {}
    for paragraph in paragraphs:
        grouped.setdefault(paragraph.slide_index, []).append(paragraph)
    max_chars, max_items = batch_limits()
    total_batches = 0
    for slide_paragraphs in grouped.values():
        segments = []
        for idx, paragraph in enumerate(slide_paragraphs):
            translatable_cores = [
                core for core, can_translate in zip(paragraph.cores, paragraph.translate_mask) if can_translate
            ]
            if not translatable_cores:
                continue
            segments.append(
                TranslationSegment(segment_id=str(idx), text=build_segmented_text(translatable_cores))
            )
        total_batches += sum(1 for _ in batch_segments(segments, max_chars=max_chars, max_items=max_items))
    return total_batches


def build_segmented_text(parts: Sequence[str]) -> str:
    chunks = []
    for idx, part in enumerate(parts):
        chunks.append(f"[[[RUN_{idx}]]]")
        chunks.append(encode_line_breaks(part))
    chunks.append("[[[RUN_END]]]")
    return "".join(chunks)


def split_segmented_text(text: str, count: int) -> List[str]:
    pattern = re.compile(r"\[\[\[RUN_(\d+|END)\]\]\]")
    matches = list(pattern.finditer(text))
    if not matches:
        raise ValueError("No run markers found in translation.")
    segments = {}
    for index, match in enumerate(matches):
        token = match.group(1)
        start = match.end()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        if token == "END":
            continue
        segments[int(token)] = text[start:end]
    if len(segments) != count:
        raise ValueError("Run marker count mismatch after translation.")
    return [decode_line_breaks(segments[i]) for i in range(count)]


@dataclass
class ProgressTracker:
    total: int
    container: object
    completed: int = 0

    def advance(self) -> None:
        self.completed += 1
        fraction = self.completed / self.total if self.total else 1.0
        render_progress(self.container, fraction)


def translate_paragraphs(
    paragraphs: List[ParagraphInfo],
    config: TranslationConfig,
    tracker: ProgressTracker | None = None,
    lexicon: LexiconStore | None = None,
) -> None:
    grouped: Dict[int, List[ParagraphInfo]] = {}
    for paragraph in paragraphs:
        grouped.setdefault(paragraph.slide_index, []).append(paragraph)

    max_chars, max_items = batch_limits()

    for slide_index in sorted(grouped.keys()):
        slide_paragraphs = grouped[slide_index]
        paragraph_maps: List[List[int | None]] = []
        segments: List[TranslationSegment] = []
        segment_matches: Dict[str, List[LexiconEntry]] = {}
        segment_texts: Dict[str, str] = {}

        for idx, paragraph in enumerate(slide_paragraphs):
            run_map: List[int | None] = []
            translatable_cores: List[str] = []
            for core, can_translate in zip(paragraph.cores, paragraph.translate_mask):
                if can_translate:
                    run_map.append(len(translatable_cores))
                    translatable_cores.append(core)
                else:
                    run_map.append(None)
            segmented = build_segmented_text(translatable_cores)
            segment_id = str(idx)
            segments.append(TranslationSegment(segment_id=segment_id, text=segmented))
            if lexicon:
                segment_matches[segment_id] = select_lexicon_matches(segmented, lexicon)
            segment_texts[segment_id] = segmented
            paragraph_maps.append(run_map)

        translated_map: Dict[str, str] = {}
        for batch in batch_segments(segments, max_chars=max_chars, max_items=max_items):
            batch_key = tuple((seg.segment_id, seg.text) for seg in batch)
            rag_context = (
                format_lexicon_context(
                    merge_lexicon_matches(segment_matches.get(seg.segment_id, []) for seg in batch)
                )
                if lexicon
                else None
            )
            try:
                translated_map.update(translate_segments_cached(batch_key, config, rag_context))
            except Exception:
                for seg in batch:
                    seg_context = (
                        format_lexicon_context(segment_matches.get(seg.segment_id, [])) if lexicon else None
                    )
                    translated_map[seg.segment_id] = cached_translate(seg.text, config, seg_context)
            for seg in batch:
                if seg.segment_id in translated_map:
                    translated_map[seg.segment_id] = postprocess_translation(
                        seg.text,
                        translated_map[seg.segment_id],
                        config,
                        segment_matches.get(seg.segment_id),
                    )

        for idx, paragraph in enumerate(slide_paragraphs):
            translated = translated_map.get(str(idx), "")
            run_map = paragraph_maps[idx]
            translatable_count = sum(1 for item in run_map if item is not None)
            try:
                translated_parts = split_segmented_text(translated, translatable_count)
            except Exception:
                original_segment = segment_texts.get(str(idx), translated)
                repaired = repair_translation(
                    original_segment,
                    translated,
                    config,
                    segment_matches.get(str(idx), []),
                )
                try:
                    translated_parts = split_segmented_text(repaired, translatable_count)
                except Exception:
                    translated_parts = []
                    for core, can_translate in zip(paragraph.cores, paragraph.translate_mask):
                        if can_translate:
                            glossary = select_lexicon_matches(core, lexicon) if lexicon else None
                            translated_parts.append(translate_text_with_postprocess(core, config, glossary))

            translated_length = 0
            for run, prefix, suffix, core, map_index, can_translate in zip(
                paragraph.runs,
                paragraph.prefixes,
                paragraph.suffixes,
                paragraph.cores,
                run_map,
                paragraph.translate_mask,
            ):
                if not can_translate:
                    run.text = f"{prefix}{core}{suffix}"
                    continue
                if map_index is not None and map_index < len(translated_parts):
                    translated_core = translated_parts[map_index]
                else:
                    translated_core = core
                translated_length += len(translated_core)
                run.text = f"{prefix}{translated_core}{suffix}"

            if not paragraph.is_table and paragraph.shape is not None:
                text_frame = paragraph.shape.text_frame
                if paragraph.is_title:
                    text_frame.word_wrap = False
                else:
                    text_frame.word_wrap = True
                if paragraph.original_length and translated_length > paragraph.original_length * 1.05:
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if tracker:
                tracker.advance()


def translate_xml_text_nodes(
    xml_bytes: bytes,
    config: TranslationConfig,
    lexicon: LexiconStore | None = None,
) -> bytes:
    tree = ET.ElementTree(ET.fromstring(xml_bytes))
    root = tree.getroot()
    text_nodes = []
    for elem in root.iter():
        if elem.tag.endswith("}t"):
            text = elem.text or ""
            if not text.strip():
                continue
            if is_numeric_text(text) or is_math_text(text):
                continue
            text_nodes.append(elem)

    if not text_nodes:
        return xml_bytes

    segments = [
        TranslationSegment(segment_id=str(i), text=encode_line_breaks(node.text or ""))
        for i, node in enumerate(text_nodes)
    ]
    segment_matches: Dict[str, List[LexiconEntry]] = {}
    if lexicon:
        for seg in segments:
            segment_matches[seg.segment_id] = select_lexicon_matches(seg.text, lexicon)
    translated_map: Dict[str, str] = {}
    max_chars, max_items = batch_limits()
    for batch in batch_segments(segments, max_chars=max_chars, max_items=max_items):
        batch_key = tuple((seg.segment_id, seg.text) for seg in batch)
        rag_context = (
            format_lexicon_context(merge_lexicon_matches(segment_matches.get(seg.segment_id, []) for seg in batch))
            if lexicon
            else None
        )
        try:
            translated_map.update(translate_segments_cached(batch_key, config, rag_context))
        except Exception:
            for seg in batch:
                seg_context = (
                    format_lexicon_context(segment_matches.get(seg.segment_id, [])) if lexicon else None
                )
                translated_map[seg.segment_id] = cached_translate(seg.text, config, seg_context)
        for seg in batch:
            if seg.segment_id in translated_map:
                translated_map[seg.segment_id] = postprocess_translation(
                    seg.text,
                    translated_map[seg.segment_id],
                    config,
                    segment_matches.get(seg.segment_id),
                )

    for idx, node in enumerate(text_nodes):
        translated = translated_map.get(str(idx))
        if translated:
            node.text = decode_line_breaks(translated)
    output = BytesIO()
    tree.write(output, encoding="utf-8", xml_declaration=True)
    return output.getvalue()


def check_ocr_dependencies() -> str | None:
    try:
        import pytesseract  # noqa: F401
        from PIL import Image  # noqa: F401
    except Exception:
        return "OCR dependencies are missing. Install pillow and pytesseract, and ensure Tesseract is installed."
    return None


def extract_text_from_image_ocr(image_bytes: bytes) -> str:
    if not OCR_IMAGE_ENABLED:
        return ""
    try:
        from PIL import Image, ImageOps
        import pytesseract
    except Exception:
        return ""

    try:
        image = Image.open(BytesIO(image_bytes))
    except Exception:
        return ""
    try:
        image = ImageOps.exif_transpose(image)
        image = ImageOps.grayscale(image)
    except Exception:
        pass
    text = pytesseract.image_to_string(image)
    return text.strip()


def replace_picture_with_textbox(slide, shape, text: str) -> None:
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    parent = shape.element.getparent()
    if parent is not None:
        parent.remove(shape.element)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = "IMAGE_TEXT_TRANSLATED"
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = text
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def count_picture_shapes(presentation: Presentation) -> int:
    count = 0
    for slide in presentation.slides:
        for shape in iter_slide_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                count += 1
    return count


def translate_image_shapes(
    presentation: Presentation,
    config: TranslationConfig,
    lexicon: LexiconStore | None = None,
    tracker: ProgressTracker | None = None,
) -> None:
    if not OCR_IMAGE_ENABLED:
        return
    for slide in presentation.slides:
        shapes = [shape for shape in iter_slide_shapes(slide.shapes) if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        for shape in shapes:
            extracted = extract_text_from_image_ocr(shape.image.blob)
            if extracted.strip():
                glossary = select_lexicon_matches(extracted, lexicon) if lexicon else None
                translated = translate_text_with_postprocess(extracted, config, glossary)
                replace_picture_with_textbox(slide, shape, translated)
            if tracker:
                tracker.advance()


def count_xml_parts(pptx_bytes: bytes) -> int:
    with zipfile.ZipFile(BytesIO(pptx_bytes), "r") as zin:
        return sum(
            1
            for item in zin.infolist()
            if item.filename.endswith(".xml")
            and (
                item.filename.startswith("ppt/charts/")
                or item.filename.startswith("ppt/diagrams/")
                or item.filename.startswith("ppt/diagramData/")
                or item.filename.startswith("ppt/diagramLayout/")
                or item.filename.startswith("ppt/diagramStyles/")
            )
        )


def translate_xml_parts(
    pptx_bytes: bytes,
    config: TranslationConfig,
    tracker: ProgressTracker | None = None,
    lexicon: LexiconStore | None = None,
) -> bytes:
    input_buffer = BytesIO(pptx_bytes)
    output_buffer = BytesIO()
    with zipfile.ZipFile(input_buffer, "r") as zin, zipfile.ZipFile(
        output_buffer, "w", compression=zipfile.ZIP_DEFLATED
    ) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/charts/")
                or item.filename.startswith("ppt/diagrams/")
                or item.filename.startswith("ppt/diagramData/")
                or item.filename.startswith("ppt/diagramLayout/")
                or item.filename.startswith("ppt/diagramStyles/")
            ):
                try:
                    data = translate_xml_text_nodes(data, config, lexicon)
                except Exception:
                    pass
                if tracker:
                    tracker.advance()
            zout.writestr(item, data)
    output_buffer.seek(0)
    return output_buffer.read()


def set_status(container, message: str) -> None:
    container.markdown(
        "<div style='border:1px solid #a8a29f; border-left:4px solid #bc0031; "
        "padding:12px; background:#ffffff; color:#1B1918;'>"
        f"{message}"
        "</div>",
        unsafe_allow_html=True,
    )


def render_progress(container, fraction: float) -> None:
    percent = max(0.0, min(1.0, fraction)) * 100
    container.markdown(
        f"<div class='progress-track'><div class='progress-fill' style='width:{percent:.2f}%;'></div></div>",
        unsafe_allow_html=True,
    )


def main() -> None:
    st.set_page_config(page_title="PPTX Translator", page_icon="🈯", layout="wide")
    st.title("PPTX Translator")
    st.markdown(
        """
        <style>
        .stApp { background-color: #ffffff; color: #1B1918; }
        .block-container { padding-top: 2rem; max-width: 100% !important; padding-left: 2rem; padding-right: 2rem; }
        [data-testid="stSidebar"] { background-color: #f2f2f2; }
        h1, h2, h3, h4, h5, h6 { color: #1B1918; }
        * { border-radius: 0 !important; }
        .stButton > button {
            background-color: #bc0031;
            color: #ffffff;
            border-radius: 0;
            border: none;
        }
        .stButton > button:hover { background-color: #a1002a; color: #ffffff; }
        .stTextInput > div > div > input,
        .stFileUploader,
        .stProgress > div > div,
        [data-testid="stAlert"],
        .stTextInput,
        .stSelectbox,
        .stRadio {
            border-radius: 0 !important;
        }
        [data-testid="stAlert"] {
            background: #ffffff !important;
            border: 1px solid #a8a29f !important;
            border-left: 4px solid #bc0031 !important;
            color: #1B1918 !important;
        }
        [data-testid="stAlert"] svg { color: #bc0031 !important; }
        [data-testid="stAlert"] [data-testid="stMarkdownContainer"] { color: #1B1918 !important; }
        .progress-track {
            width: 100%;
            height: 14px;
            background: #e5e5e5;
            border: 1px solid #a8a29f;
        }
        .progress-fill {
            height: 100%;
            background: #bc0031;
            width: 0%;
        }
        div[data-testid="stNotification"] { display: none; }
        </style>
        """,
        unsafe_allow_html=True,
    )
    st.markdown(
        "Upload a PowerPoint deck, translate the text slide-by-slide, "
        "and download a PPTX with the original layout preserved."
    )

    with st.sidebar:
        st.subheader("Translation Settings")
        target_language = st.text_input("Target language", value="Dutch")
        api_key = st.text_input("API key", type="password")
        st.caption("Uses LiteLLM endpoint with model gpt-oss-120b.")

    uploaded_file = st.file_uploader("Drop your PPTX or DOCX file here", type=["pptx", "docx"])
    if not uploaded_file:
        st.markdown(
            "<div style='border:1px solid #a8a29f; border-left:4px solid #bc0031; "
            "padding:12px; background:#ffffff; color:#1B1918;'>"
            "Upload a .pptx or .docx file to get started."
            "</div>",
            unsafe_allow_html=True,
        )
        return

    lexicon_store = None
    if target_language.strip().lower() in {"dutch", "nederlands", "nl"}:
        lexicon_store = load_finance_lexicon()
        if not lexicon_store:
            set_status(
                st.empty(),
                f"Finance lexicon PDF not found or unreadable: {LEXICON_PDF_NAME}.",
            )
            return
    config = TranslationConfig(
        base_url=DEFAULT_BASE_URL,
        model=DEFAULT_MODEL,
        target_language=target_language,
        temperature=DEFAULT_TEMPERATURE,
        timeout=DEFAULT_TIMEOUT,
        api_key=api_key,
    )

    file_id = f"{uploaded_file.name}:{uploaded_file.size}"
    if st.session_state.get("file_id") != file_id:
        st.session_state["file_id"] = file_id
        st.session_state.pop("translated_file", None)
        st.session_state.pop("translated_name", None)
        st.session_state.pop("translated_mime", None)
        st.session_state.pop("translated_label", None)

    action_col, download_col, _spacer = st.columns([1, 1, 6], gap="small")
    translate_clicked = action_col.button("Translate Deck", type="primary")
    download_placeholder = download_col.empty()
    if "translated_file" in st.session_state:
        download_placeholder.download_button(
            st.session_state.get("translated_label", "Download File"),
            data=st.session_state["translated_file"],
            file_name=st.session_state.get("translated_name", uploaded_file.name),
            mime=st.session_state.get("translated_mime", "application/octet-stream"),
        )
    if translate_clicked:
        if not config.api_key:
            set_status(st.empty(), "Please provide an API key.")
            return
        raw_bytes = uploaded_file.read()
        if not raw_bytes:
            set_status(st.empty(), "The uploaded file was empty.")
            return

        file_suffix = Path(uploaded_file.name).suffix.lower()
        progress_container = st.empty()
        status_text = st.empty()

        if file_suffix == ".pptx":
            presentation = Presentation(BytesIO(raw_bytes))
            paragraphs = collect_paragraphs(presentation)

            if not paragraphs:
                set_status(st.empty(), "No translatable paragraphs were found in this deck.")
                return

            xml_count = count_xml_parts(raw_bytes)
            raw_image_count = count_picture_shapes(presentation)
            ocr_error = None
            if OCR_IMAGE_ENABLED and raw_image_count:
                ocr_error = check_ocr_dependencies()
                if ocr_error:
                    set_status(st.empty(), ocr_error)
                    return
            image_count = raw_image_count if OCR_IMAGE_ENABLED else 0
            total_steps = len(paragraphs) + xml_count + image_count
            if total_steps == 0:
                total_steps = 1

            render_progress(progress_container, 0.0)
            tracker = ProgressTracker(total=total_steps, container=progress_container)
            start_time = time.time()

            try:
                set_status(status_text, "Translating text…")
                translate_paragraphs(paragraphs, config, tracker=tracker, lexicon=lexicon_store)
            except requests.RequestException as exc:
                set_status(st.empty(), f"Translation request failed: {exc}")
                return
            except RuntimeError as exc:
                set_status(st.empty(), str(exc))
                return
            except Exception as exc:
                set_status(st.empty(), f"Translation failed: {exc}")
                return

            if image_count:
                set_status(status_text, "Translating image text…")
                translate_image_shapes(presentation, config, lexicon=lexicon_store, tracker=tracker)

            output = BytesIO()
            presentation.save(output)
            pptx_bytes = output.getvalue()
            if xml_count:
                set_status(status_text, "Translating chart/SmartArt text…")
            pptx_bytes = translate_xml_parts(pptx_bytes, config, tracker=tracker, lexicon=lexicon_store)
            output = BytesIO(pptx_bytes)
            output.seek(0)

            render_progress(progress_container, 1.0)

            elapsed = time.time() - start_time
            set_status(status_text, f"Translation complete in {elapsed:.1f}s.")

            st.session_state["translated_file"] = output.getvalue()
            st.session_state["translated_name"] = (
                f"{uploaded_file.name.rsplit('.', 1)[0]}_{target_language}.pptx"
            )
            st.session_state["translated_mime"] = (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            st.session_state["translated_label"] = "Download PPTX"
        elif file_suffix == ".docx":
            document = Document(BytesIO(raw_bytes))
            paragraphs = collect_docx_paragraphs(document)
            if not paragraphs:
                set_status(st.empty(), "No translatable paragraphs were found in this document.")
                return

            total_steps = len(paragraphs)
            if total_steps == 0:
                total_steps = 1

            render_progress(progress_container, 0.0)
            tracker = ProgressTracker(total=total_steps, container=progress_container)
            start_time = time.time()

            try:
                set_status(status_text, "Translating text…")
                translate_paragraphs(paragraphs, config, tracker=tracker, lexicon=lexicon_store)
            except requests.RequestException as exc:
                set_status(st.empty(), f"Translation request failed: {exc}")
                return
            except RuntimeError as exc:
                set_status(st.empty(), str(exc))
                return
            except Exception as exc:
                set_status(st.empty(), f"Translation failed: {exc}")
                return

            output = BytesIO()
            document.save(output)
            output.seek(0)

            render_progress(progress_container, 1.0)

            elapsed = time.time() - start_time
            set_status(status_text, f"Translation complete in {elapsed:.1f}s.")

            st.session_state["translated_file"] = output.getvalue()
            st.session_state["translated_name"] = (
                f"{uploaded_file.name.rsplit('.', 1)[0]}_{target_language}.docx"
            )
            st.session_state["translated_mime"] = (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
            st.session_state["translated_label"] = "Download DOCX"
        else:
            set_status(st.empty(), "Unsupported file type. Please upload a PPTX or DOCX file.")
            return

        download_placeholder.download_button(
            st.session_state.get("translated_label", "Download File"),
            data=st.session_state["translated_file"],
            file_name=st.session_state["translated_name"],
            mime=st.session_state["translated_mime"],
        )

if __name__ == "__main__":
    main()
